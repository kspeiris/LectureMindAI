"""
Upgraded PPTX processor: skips non-content shapes (tables, pictures),
uses the advanced text_cleaner for artifact-free output.
"""
from pptx import Presentation
from pptx.util import Pt
import os
from .text_cleaner import clean_text


def extract_text_from_ppt(filepath: str) -> dict:
    """
    Extracts text from a PPTX file using python-pptx.
    - Skips image/chart/table shapes (they produce noise)
    - Prefers slide body text over title-only slides
    - Applies advanced clean_text post-processing

    Returns:
        dict: {'text': str, 'pages': int, 'word_count': int}
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File not found: {filepath}")

    prs = Presentation(filepath)
    slide_texts = []
    num_slides = len(prs.slides)

    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            # Skip non-text shapes
            if not hasattr(shape, "text_frame"):
                continue
            # Skip very short title-only text (< 4 words)
            raw = shape.text_frame.text.strip()
            if not raw or len(raw.split()) < 3:
                continue
            # Gather paragraph text, preserving structure
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if para_text and len(para_text.split()) >= 2:
                    parts.append(para_text)

        if parts:
            slide_texts.append("\n".join(parts))

    raw_text    = "\n\n".join(slide_texts)
    cleaned     = clean_text(raw_text)
    word_count  = len(cleaned.split())

    return {
        "text":       cleaned,
        "pages":      num_slides,
        "word_count": word_count,
    }
